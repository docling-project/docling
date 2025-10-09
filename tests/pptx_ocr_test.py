from pptx import Presentation
import easyocr
from PIL import Image
import io
import numpy as np

pptx_file = r"C:\Users\Admin\Desktop\demo\docling\tests\example.pptx"
reader = easyocr.Reader(['en'])

# Load PPTX
prs = Presentation(pptx_file)

for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture
            image = shape.image
            image_bytes = io.BytesIO(image.blob)
            img = Image.open(image_bytes).convert('RGB')
            # Convert PIL Image to Numpy array
            img_np = np.array(img)
            # OCR the image
            result = reader.readtext(img_np)
            print(f"Slide {i+1}, Image OCR Results:")
            for (bbox, text, prob) in result:
                print(text)
