from pathlib import Path

import pytest
from docling_core.types.doc import DocItemLabel
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches

from docling.datamodel.backend_options import MsPowerpointBackendOptions
from docling.datamodel.base_models import InputFormat
from docling.datamodel.document import ConversionResult, DoclingDocument
from docling.document_converter import DocumentConverter, PowerpointFormatOption

from .test_data_gen_flag import GEN_TEST_DATA
from .verify_utils import verify_document, verify_export

GENERATE = GEN_TEST_DATA


def get_pptx_paths():
    # Define the directory you want to search
    directory = Path("./tests/data/pptx/")

    # List all PPTX files in the directory and its subdirectories
    pptx_files = sorted(directory.rglob("*.pptx"))
    return pptx_files


def get_converter():
    converter = DocumentConverter(allowed_formats=[InputFormat.PPTX])

    return converter


def test_e2e_pptx_conversions():
    pptx_paths = get_pptx_paths()
    converter = get_converter()

    for pptx_path in pptx_paths:
        # print(f"converting {pptx_path}")

        gt_path = (
            pptx_path.parent.parent / "groundtruth" / "docling_v2" / pptx_path.name
        )

        conv_result: ConversionResult = converter.convert(pptx_path)

        doc: DoclingDocument = conv_result.document

        pred_md: str = doc.export_to_markdown(compact_tables=True)
        assert verify_export(pred_md, str(gt_path) + ".md", GENERATE), "export to md"

        pred_itxt: str = doc._export_to_indented_text(
            max_text_len=70, explicit_tables=False
        )
        assert verify_export(pred_itxt, str(gt_path) + ".itxt", GENERATE), (
            "export to indented-text"
        )

        assert verify_document(doc, str(gt_path) + ".json", GENERATE), (
            "document document"
        )


def test_pptx_unrecognized_shape_type():
    """PPTX with a <p:sp> that has no geometry should not crash.

    python-pptx raises NotImplementedError from Shape.shape_type for shapes
    that aren't placeholders, autoshapes, textboxes, or freeforms. The
    backend should skip the unrecognized shape gracefully and still extract
    text from the rest of the presentation.

    Ref: https://github.com/docling-project/docling/issues/3308
    """
    converter = get_converter()
    pptx_path = Path("./tests/data/pptx/powerpoint_unrecognized_shape.pptx")

    conv_result: ConversionResult = converter.convert(pptx_path)
    doc: DoclingDocument = conv_result.document

    pred_md = doc.export_to_markdown()

    # Normal slide content should still be extracted
    assert "Q3 Revenue Summary" in pred_md
    assert "Enterprise segment" in pred_md
    assert "Key Metrics" in pred_md
    assert "Next Steps" in pred_md


def test_pptx_malformed_picture_shapes():
    """PPTX with malformed <p:pic> shapes should not crash conversion.

    python-pptx's shape.image accessor raises three distinct exceptions on
    picture shapes that slip past other tools' parsers (Keynote/Google Drive
    open these files fine): InvalidXmlError when <p:blipFill> is missing,
    KeyError when <a:blip r:embed> points at an unknown relationship, and
    AttributeError when the embedded part's content-type isn't an image.

    The backend should skip each malformed picture with a warning and still
    extract text from the slides.
    """
    converter = get_converter()
    pptx_path = Path("./tests/data/pptx/powerpoint_malformed_pictures.pptx")

    with pytest.warns(UserWarning, match="Skipping malformed picture shape"):
        conv_result: ConversionResult = converter.convert(pptx_path)

    doc: DoclingDocument = conv_result.document

    pred_md = doc.export_to_markdown()
    assert "Slide With Missing BlipFill" in pred_md
    assert "Slide With Dangling Rel" in pred_md
    assert "Slide With Wrong Content Type" in pred_md


def test_pptx_page_range():
    converter = get_converter()
    pptx_path = Path("./tests/data/pptx/powerpoint_sample.pptx")

    conv_result: ConversionResult = converter.convert(pptx_path, page_range=(2, 2))

    assert conv_result.input.page_count == 3
    assert conv_result.document.num_pages() == 1
    assert list(conv_result.document.pages.keys()) == [2]

    pred_md = conv_result.document.export_to_markdown()
    assert "Second slide title" in pred_md
    assert "Test Table Slide" not in pred_md
    assert "List item4" not in pred_md


def _make_pptx_with_two_images_and_connector(tmp_path):
    left_img = tmp_path / "left.png"
    right_img = tmp_path / "right.png"
    pptx_path = tmp_path / "composition_arrow.pptx"

    for img_path, label in [(left_img, "A"), (right_img, "B")]:
        image = Image.new("RGB", (320, 220), "white")
        draw = ImageDraw.Draw(image)
        draw.rectangle((20, 20, 300, 200), outline="black", width=5)
        draw.text((145, 100), label, fill="black")
        image.save(img_path)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.shapes.add_picture(
        str(left_img),
        Inches(1),
        Inches(2),
        width=Inches(2.5),
        height=Inches(1.8),
    )
    slide.shapes.add_picture(
        str(right_img),
        Inches(6),
        Inches(2),
        width=Inches(2.5),
        height=Inches(1.8),
    )
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(3.7),
        Inches(2.9),
        Inches(5.8),
        Inches(2.9),
    )

    prs.save(pptx_path)
    return pptx_path


def test_pptx_slide_visual_composition_is_opt_in(tmp_path):
    pptx_path = _make_pptx_with_two_images_and_connector(tmp_path)

    converter = DocumentConverter(allowed_formats=[InputFormat.PPTX])
    result = converter.convert(pptx_path)

    pictures = [
        item
        for item, _level in result.document.iterate_items()
        if getattr(item, "label", None) == DocItemLabel.PICTURE
    ]

    assert len(pictures) == 2


def test_pptx_slide_visual_composition_adds_slide_picture(tmp_path):
    pptx_path = _make_pptx_with_two_images_and_connector(tmp_path)

    converter = DocumentConverter(
        format_options={
            InputFormat.PPTX: PowerpointFormatOption(
                backend_options=MsPowerpointBackendOptions(
                    create_slide_visual_composition=True,
                )
            )
        }
    )
    result = converter.convert(pptx_path)

    pictures = [
        item
        for item, _level in result.document.iterate_items()
        if getattr(item, "label", None) == DocItemLabel.PICTURE
    ]

    assert len(pictures) == 3

    slide_picture = pictures[0]
    assert slide_picture.prov
    assert slide_picture.prov[0].page_no == 1
    assert slide_picture.prov[0].bbox.l == 0
    assert slide_picture.prov[0].bbox.b == 0
