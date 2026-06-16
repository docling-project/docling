import logging
from io import BytesIO

from PIL import Image, ImageDraw, UnidentifiedImageError
from pptx.enum.shapes import MSO_SHAPE_TYPE

_log = logging.getLogger(__name__)

EMU_PER_INCH = 914400


class PptxSlideCompositionRenderer:
    """Lightweight renderer for PPTX slide-level visual composition.

    This is intentionally not a complete PowerPoint renderer. It captures the
    main visual relationships needed by downstream VLM enrichment, especially
    images connected by simple line/connector shapes.
    """

    def __init__(self, dpi: int = 144) -> None:
        self.dpi = dpi

    def emu_to_px(self, value: int | float) -> int:
        return round(float(value) / EMU_PER_INCH * self.dpi)

    def render(
        self,
        slide,
        slide_width: int | float,
        slide_height: int | float,
    ) -> Image.Image:
        width_px = max(1, self.emu_to_px(slide_width))
        height_px = max(1, self.emu_to_px(slide_height))

        canvas = Image.new("RGB", (width_px, height_px), "white")
        draw = ImageDraw.Draw(canvas)

        for shape in slide.shapes:
            self._render_shape(shape, canvas, draw)

        return canvas

    def _render_shape(
        self,
        shape,
        canvas: Image.Image,
        draw: ImageDraw.ImageDraw,
    ) -> None:
        shape_type = self._safe_shape_type(shape)

        if shape_type is None:
            return

        if shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                self._render_shape(child, canvas, draw)
            return

        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            self._render_picture(shape, canvas)
            return

        if shape_type == MSO_SHAPE_TYPE.LINE:
            self._render_line(shape, draw)
            return

        if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            self._render_auto_shape(shape, draw)

        if getattr(shape, "has_text_frame", False):
            self._render_text(shape, draw)

    def _safe_shape_type(self, shape):
        try:
            return shape.shape_type
        except NotImplementedError:
            _log.debug("Skipping shape with unrecognized type in slide rendering")
            return None

    def _render_picture(self, shape, canvas: Image.Image) -> None:
        try:
            image = Image.open(BytesIO(shape.image.blob)).convert("RGBA")
        except (
            UnidentifiedImageError,
            OSError,
            ValueError,
            KeyError,
            AttributeError,
        ) as exc:
            _log.debug("Skipping malformed picture in slide composition: %s", exc)
            return

        x = self.emu_to_px(shape.left)
        y = self.emu_to_px(shape.top)
        w = max(1, self.emu_to_px(shape.width))
        h = max(1, self.emu_to_px(shape.height))

        image = image.resize((w, h))
        canvas.paste(image, (x, y), image)

    def _render_line(self, shape, draw: ImageDraw.ImageDraw) -> None:
        x1 = self.emu_to_px(shape.left)
        y1 = self.emu_to_px(shape.top)
        x2 = self.emu_to_px(shape.left + shape.width)
        y2 = self.emu_to_px(shape.top + shape.height)

        if x1 == x2 and y1 == y2:
            return

        draw.line(
            (x1, y1, x2, y2),
            fill="black",
            width=self._line_width_px(shape),
        )

    def _render_auto_shape(self, shape, draw: ImageDraw.ImageDraw) -> None:
        x1 = self.emu_to_px(shape.left)
        y1 = self.emu_to_px(shape.top)
        x2 = self.emu_to_px(shape.left + shape.width)
        y2 = self.emu_to_px(shape.top + shape.height)

        if x1 == x2 or y1 == y2:
            return

        draw.rectangle(
            (x1, y1, x2, y2),
            outline="black",
            width=self._line_width_px(shape),
        )

    def _render_text(self, shape, draw: ImageDraw.ImageDraw) -> None:
        text = getattr(shape, "text", "") or ""
        text = text.strip()

        if not text:
            return

        x = self.emu_to_px(shape.left) + 4
        y = self.emu_to_px(shape.top) + 4

        draw.text((x, y), text, fill="black")

    def _line_width_px(self, shape) -> int:
        try:
            width = shape.line.width
        except AttributeError:
            width = None

        if width is None:
            return 2

        return max(1, self.emu_to_px(width))
